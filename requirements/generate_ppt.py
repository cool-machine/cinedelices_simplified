from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()

    # Helper to add standard slide with title and text
    def add_slide(title, content_lines, image_path=None):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Content - Use a text box for more control if image is present, 
        # or defaults. Here we stick to placeholder but adjust font.
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        # Reduce text size to prevent overflow
        for i, line in enumerate(content_lines):
            if i == 0:
                tf.text = line
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.text = line
            
            p.font.size = Pt(18) # Slightly larger than 14 for readability, still safe
            
            # Formatting
            if line.startswith('* '):
                 p.text = line[2:]
                 p.level = 0
            elif line.startswith('    * '):
                 p.text = line[6:]
                 p.level = 1

        # Add Image if provided
        if image_path and os.path.exists(image_path):
            # Add image on the right or bottom
            # For simplicity, we put it in a fixed position (bottom right usually)
            # or we resize the text box.
            
            # Move text box to left half
            body_shape.width = Inches(4.5)
            
            # Add image in right half
            left = Inches(5)
            top = Inches(1.5)
            height = Inches(4.5)
            # We let width be auto to preserve aspect ratio, or constrain both
            slide.shapes.add_picture(image_path, left, top, height=height)

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = "Présentation Sprint 0 - Conception"
    subtitle.text = "Projet : Ciné Délices\nBilan de la phase de conception"

    # 2. Intro
    add_slide("1. Introduction", [
        "Le Projet Ciné Délices 🎬 🍽️",
        "",
        "Concept : Mariage du Cinéma et de la Gastronomie.",
        "Objectif : Reproduire les plats iconiques des films.",
        "",
        "L'Équipe :",
        "* Lead Backend : Architecture & API",
        "* Lead Frontend : UX & Intégration",
        "* DevOps : CI/CD & Déploiement",
        "* Design : Identité & Maquettes"
    ])

    # 3. Specs
    add_slide("2. Cahier des Charges", [
        "Vision & MVP :",
        "* 🔐 Authentification complète",
        "* 📖 Catalogue de recettes (filtres/recherche)",
        "* 📄 Fiches Recettes (ingrédients, films associés)",
        "* ⚙️ Back-office Admin",
        "",
        "Cible :",
        "* Passionnés de pop-culture",
        "* Cuisiniers amateurs"
    ])

    # 4. Architecture (with Image)
    add_slide("3. Architecture", [
        "Arborescence du site :",
        "",
        "* Accueil",
        "* Catalogue Recettes / Films",
        "* Espace Connecté (Profil, Favoris)",
        "* Administration",
        "",
        "Une structure optimisée pour la navigation."
    ], image_path="requirements/presentation_assets/arborescence.png")

    # 5. Design (with Image)
    add_slide("4. Design & UX", [
        "Identité Visuelle :",
        "* Ambiance 'Salle Obscure' & 'Tapis Rouge'",
        "* Couleurs : Or (#D4AF37), Rouge (#8B0000), Noir",
        "* Typo : Bebas Neue (Titres), Montserrat (Corps)",
        "",
        "Maquette Accueil (Mode Sombre) ->"
    ], image_path="requirements/presentation_assets/homepage.png")

    # 6. Mockups Recette (with Image)
    add_slide("4.1 Maquette Recette", [
        "Détail de la Recette :",
        "* Mise en page immersive",
        "* Bloc Ingrédients clair",
        "* Instructions étape par étape",
        "* Rappel du film source",
        "",
        "Focus sur l'expérience utilisateur."
    ], image_path="requirements/presentation_assets/recipe.png")

    # 7. BDD (with Image)
    add_slide("5. Base de Données", [
        "Modélisation (MCD/MPD) :",
        "* 4 Entités principales (MVP)",
        "* Relations optimisées",
        "* Intégrité référentielle forte",
        "",
        "Schéma Conceptuel ->"
    ], image_path="requirements/presentation_assets/mcd.png")

    # 8. Conclusion
    add_slide("6. Conclusion", [
        "Bilan Sprint 0 :",
        "* ✅ Conception validée",
        "* ✅ Maquettes prêtes",
        "* ✅ Base de données définie",
        "",
        "Sprint 1 :",
        "* Initialisation technique",
        "* Création API & BDD",
        "* Démarrage Frontend"
    ])

    prs.save('requirements/presentation_sprint_0.pptx')
    print("Présentation générée avec images : requirements/presentation_sprint_0.pptx")

if __name__ == "__main__":
    create_presentation()
