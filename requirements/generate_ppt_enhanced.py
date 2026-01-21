from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand Colors
COLOR_GOLD = RGBColor(212, 175, 55)      # #D4AF37
COLOR_RED = RGBColor(139, 0, 0)          # #8B0000
COLOR_BLACK = RGBColor(26, 26, 26)       # #1A1A1A
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_DARK_GRAY = RGBColor(60, 60, 60)

def create_enhanced_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    def apply_brand_colors(shape):
        """Apply brand colors to text"""
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    run.font.color.rgb = COLOR_BLACK

    def add_title_slide():
        """Create branded title slide"""
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.shapes.placeholders[1]
        
        title.text = "Cin√© D√©lices"
        title.text_frame.paragraphs[0].font.size = Pt(48)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = COLOR_GOLD
        
        subtitle.text = "Pr√©sentation Sprint 0 - Phase de Conception\n\nüé¨ Mariage du Cin√©ma et de la Gastronomie üçΩÔ∏è"
        subtitle.text_frame.paragraphs[0].font.size = Pt(18)
        subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_GRAY
        
        # Add decorative background shape
        left = Inches(0)
        top = Inches(0)
        width = Inches(10)
        height = Inches(1.2)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_RED
        shape.line.fill.background()
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.insert(2, shape._element)

    def add_slide_with_content(title, content_lines, image_path=None, layout_type="content"):
        """Enhanced slide creation with better formatting and proper margins"""
        if layout_type == "blank":
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title manually with safe margins
            left = Inches(0.5)
            top = Inches(0.3)
            width = Inches(9)
            height = Inches(0.7)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = COLOR_RED
            title_frame.margin_left = Inches(0)
            title_frame.margin_right = Inches(0)
            title_frame.margin_top = Inches(0)
            title_frame.margin_bottom = Inches(0)
        else:
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            title_shape.text = title
            title_shape.text_frame.paragraphs[0].font.size = Pt(32)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_RED
        
        # Content area with safe margins
        if layout_type == "blank":
            content_left = Inches(0.5)
            content_top = Inches(1.2)
            content_width = Inches(4.2) if image_path else Inches(9)
            content_height = Inches(5.8)
            body_shape = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            tf = body_shape.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.1)
            tf.margin_bottom = Inches(0.1)
        else:
            body_shape = slide.shapes.placeholders[1]
            if image_path:
                body_shape.left = Inches(0.5)
                body_shape.top = Inches(1.2)
                body_shape.width = Inches(4.2)
                body_shape.height = Inches(5.8)
            else:
                body_shape.left = Inches(0.5)
                body_shape.top = Inches(1.2)
                body_shape.width = Inches(9)
                body_shape.height = Inches(5.8)
            tf = body_shape.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.1)
            tf.margin_bottom = Inches(0.1)
        
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.auto_size = None  # Disable auto-size to control margins
        
        for i, line in enumerate(content_lines):
            if i == 0:
                tf.text = line
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.text = line
            
            # Formatting based on content with smaller fonts
            if line.startswith('* '):
                p.text = line[2:]
                p.level = 0
                p.font.size = Pt(16)
                p.font.bold = False
            elif line.startswith('    * ') or line.startswith('  * '):
                p.text = line.lstrip()[2:]
                p.level = 1
                p.font.size = Pt(14)
                p.font.bold = False
            elif line.startswith('**') and line.endswith('**'):
                p.text = line.strip('*')
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = COLOR_RED
            elif line.startswith('###'):
                p.text = line[3:].strip()
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = COLOR_RED
            elif line.strip() == '':
                p.text = ''
                p.font.size = Pt(10)
            else:
                p.font.size = Pt(16)
            
            p.font.color.rgb = COLOR_BLACK
            p.space_after = Pt(4)
            p.space_before = Pt(0)
        
        # Add image if provided with safe margins
        if image_path and os.path.exists(image_path):
            img_left = Inches(5.2)
            img_top = Inches(1.2)
            img_height = Inches(5.8)
            try:
                slide.shapes.add_picture(image_path, img_left, img_top, height=img_height)
            except Exception as e:
                print(f"Warning: Could not add image {image_path}: {e}")

    def add_two_column_slide(title, left_content, right_content):
        """Create a slide with two columns and proper margins"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title with safe margins
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = COLOR_RED
        title_frame.margin_left = Inches(0)
        title_frame.margin_right = Inches(0)
        title_frame.margin_top = Inches(0)
        title_frame.margin_bottom = Inches(0)
        
        # Left column with safe margins
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(5.8))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        left_frame.margin_left = Inches(0.1)
        left_frame.margin_right = Inches(0.1)
        left_frame.margin_top = Inches(0.1)
        left_frame.margin_bottom = Inches(0.1)
        for i, line in enumerate(left_content):
            if i == 0:
                left_frame.text = line
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
                p.text = line
            if line.startswith('**'):
                p.text = line.strip('*')
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = COLOR_RED
            else:
                p.font.size = Pt(16)
                p.font.color.rgb = COLOR_BLACK
            if line.startswith('* '):
                p.text = line[2:]
                p.level = 0
            p.space_after = Pt(4)
        
        # Right column with safe margins
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.2), Inches(5.8))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        right_frame.margin_left = Inches(0.1)
        right_frame.margin_right = Inches(0.1)
        right_frame.margin_top = Inches(0.1)
        right_frame.margin_bottom = Inches(0.1)
        for i, line in enumerate(right_content):
            if i == 0:
                right_frame.text = line
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
                p.text = line
            if line.startswith('**'):
                p.text = line.strip('*')
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = COLOR_RED
            else:
                p.font.size = Pt(16)
                p.font.color.rgb = COLOR_BLACK
            if line.startswith('* '):
                p.text = line[2:]
            p.space_after = Pt(4)

    # ========== SLIDES ==========
    
    # 1. Title Slide
    add_title_slide()
    
    # 2. Table of Contents
    add_slide_with_content("Plan de Pr√©sentation", [
        "**1. Introduction & Contexte**",
        "",
        "**2. Cahier des Charges**",
        "   * Vision & Objectifs",
        "   * Fonctionnalit√©s MVP",
        "   * Public cible",
        "",
        "**3. Architecture & Arborescence**",
        "",
        "**4. Wireframes & Navigation**",
        "",
        "**5. Design System & Identit√© Visuelle**",
        "",
        "**6. Maquettes Haute Fid√©lit√©**",
        "",
        "**7. Base de Donn√©es**",
        "   * Mod√®le Conceptuel (MCD)",
        "   * Mod√®le Logique (MLD)",
        "   * Mod√®le Physique (MPD)",
        "",
        "**8. Conclusion & Prochaines √âtapes**"
    ])
    
    # 3. Introduction
    add_slide_with_content("1. Introduction", [
        "**Le Projet Cin√© D√©lices** üé¨üçΩÔ∏è",
        "",
        "**Concept :**",
        "Une plateforme web innovante qui marie le Cin√©ma et la Gastronomie",
        "",
        "**Objectif :**",
        "Permettre aux fans de reproduire les plats iconiques de leurs films et s√©ries pr√©f√©r√©s",
        "",
        "**Valeur Ajout√©e :**",
        "* Exp√©rience immersive combinant culture pop et cuisine",
        "* Communaut√© de cin√©philes gourmands",
        "* Recettes authentiques li√©es aux ≈ìuvres cin√©matographiques"
    ])
    
    # 4. √âquipe
    add_two_column_slide("L'√âquipe", [
        "**Lead Backend**",
        "Architecture & API",
        "",
        "**Lead Frontend**",
        "UX & Int√©gration"
    ], [
        "**DevOps**",
        "CI/CD & D√©ploiement",
        "",
        "**Design/UX**",
        "Identit√© & Maquettes"
    ])
    
    # 5. Cahier des Charges - Vision
    add_slide_with_content("2. Cahier des Charges - Vision", [
        "**Vision du Projet**",
        "",
        "Cr√©er une plateforme compl√®te permettant de :",
        "",
        "* D√©couvrir des recettes inspir√©es du cin√©ma",
        "* Explorer le catalogue par film ou par recette",
        "* Suivre des instructions d√©taill√©es √©tape par √©tape",
        "* Partager ses cr√©ations culinaires",
        "* G√©rer un espace personnel (favoris, profil)",
        "",
        "**Approche MVP** :",
        "Focus sur les fonctionnalit√©s essentielles pour valider le concept"
    ])
    
    # 6. Fonctionnalit√©s MVP
    add_slide_with_content("2.1 Fonctionnalit√©s MVP", [
        "**Fonctionnalit√©s Principales**",
        "",
        "* üîê **Authentification s√©curis√©e**",
        "   * Inscription/Connexion",
        "   * Gestion de profil",
        "",
        "* üìñ **Catalogue de recettes**",
        "   * Recherche et filtres",
        "   * Navigation par film",
        "",
        "* üìÑ **Fiches recettes**",
        "   * Ingr√©dients complets",
        "   * Instructions d√©taill√©es",
        "",
        "* ‚öôÔ∏è **Back-office Admin**",
        "   * Gestion du contenu"
    ])
    
    # 7. Public Cible
    add_slide_with_content("2.2 Public Cible", [
        "**Profil des Utilisateurs**",
        "",
        "* **Amateurs de pop-culture** (16-50+ ans)",
        "   * Cin√©philes passionn√©s",
        "   * Fans de s√©ries TV",
        "",
        "* **Cuisiniers amateurs**",
        "   * D√©butants cherchant des recettes guid√©es",
        "   * Exp√©riment√©s en qu√™te d'inspiration",
        "",
        "* **Communaut√©s en ligne**",
        "   * Partage sur r√©seaux sociaux",
        "   * Cr√©ation de contenu culinaire"
    ])
    
    # 8. Architecture
    add_slide_with_content("3. Architecture & Arborescence", [
        "**Structure du Site**",
        "",
        "Navigation intuitive organis√©e en 4 sections principales :",
        "",
        "* **Accueil**",
        "   * Point d'entr√©e central",
        "   * Mise en avant des recettes",
        "",
        "* **Catalogue**",
        "   * Double entr√©e : Recettes / Films",
        "   * Filtres et recherche",
        "",
        "* **Espace Membre**",
        "   * Profil utilisateur",
        "   * Favoris",
        "   * Ajout de recettes",
        "",
        "* **Administration**",
        "   * Gestion globale du contenu"
    ], image_path="requirements/presentation_assets/arborescence.png")
    
    # 9. Wireframes
    add_slide_with_content("4. Wireframes & Navigation", [
        "**Phase de Wireframing**",
        "",
        "**Objectifs :**",
        "* D√©finir la structure de chaque page",
        "* Optimiser l'exp√©rience utilisateur",
        "* Valider les parcours de navigation",
        "",
        "**Pages Wirefram√©es :**",
        "* Page d'accueil",
        "* Catalogue recettes",
        "* Fiche recette d√©taill√©e",
        "* Page de connexion/inscription",
        "* Espace membre",
        "* Back-office admin",
        "",
        "*(Voir les wireframes d√©taill√©s dans le dossier 03)*"
    ])
    
    # 10. Design System - Couleurs
    add_slide_with_content("5. Design System - Identit√© Visuelle", [
        "**Palette de Couleurs**",
        "",
        "Ambiance \"Salle Obscure\" √©l√©gante :",
        "",
        "* **Or Cin√©ma** (#D4AF37)",
        "   * Titres et √©l√©ments premium",
        "   * Mise en valeur",
        "",
        "* **Rouge Tapis** (#8B0000)",
        "   * Accents et call-to-action",
        "   * √âl√©ments importants",
        "",
        "* **Noir Profond** (#1A1A1A)",
        "   * Fond principal (mode sombre)",
        "   * Contraste optimal",
        "",
        "**Philosophie :**",
        "√âl√©gance cin√©matographique moderne"
    ])
    
    # 11. Design System - Typographie
    add_slide_with_content("5.1 Design System - Typographie", [
        "**Choix Typographiques**",
        "",
        "* **Bebas Neue** (Titres)",
        "   * Impact visuel fort",
        "   * √âvoque les affiches de cin√©ma",
        "   * Caract√®res majuscules",
        "",
        "* **Montserrat** (Corps de texte)",
        "   * Excellente lisibilit√©",
        "   * Style moderne et √©pur√©",
        "   * Parfaite pour le web",
        "",
        "**Hi√©rarchie :**",
        "* Titres H1 : Bebas Neue, 48-54pt",
        "* Titres H2 : Bebas Neue, 36-42pt",
        "* Corps : Montserrat, 16-18pt"
    ])
    
    # 12. Maquette Accueil
    add_slide_with_content("6. Maquettes - Page d'Accueil", [
        "**Design de la Page d'Accueil**",
        "",
        "**Caract√©ristiques :**",
        "* Immersion imm√©diate",
        "* Visuels style \"Affiche de film\"",
        "* Mise en avant des recettes populaires",
        "* Navigation intuitive",
        "* Mode sombre pour l'ambiance",
        "",
        "**√âl√©ments cl√©s :**",
        "* Hero section impactante",
        "* Grille de recettes en vedette",
        "* Call-to-action clairs",
        "* Footer informatif"
    ], image_path="requirements/presentation_assets/homepage.png")
    
    # 13. Maquette Recette
    add_slide_with_content("6.1 Maquettes - Fiche Recette", [
        "**Page D√©tail Recette**",
        "",
        "**Structure :**",
        "* Image principale du plat",
        "* Informations du film source",
        "* Bloc ingr√©dients clair et lisible",
        "* Instructions √©tape par √©tape",
        "* Section commentaires",
        "* Actions (favoris, partage)",
        "",
        "**Focus UX :**",
        "* Mise en page √©ditoriale",
        "* Lisibilit√© optimale",
        "* Exp√©rience immersive"
    ], image_path="requirements/presentation_assets/recipe.png")
    
    # 14. Base de Donn√©es - MCD
    add_slide_with_content("7.1 Base de Donn√©es - MCD", [
        "**Mod√®le Conceptuel de Donn√©es (MCD)**",
        "",
        "**Entit√©s principales :**",
        "* UTILISATEUR",
        "* RECETTE",
        "* FILM",
        "* CATEGORIE",
        "",
        "**Relations :**",
        "* UTILISATEUR ‚Üî RECETTE (cr√©ation, favoris)",
        "* RECETTE ‚Üî FILM (inspiration)",
        "* RECETTE ‚Üî CATEGORIE (classification)",
        "",
        "**R√®gles m√©tier :**",
        "* Une recette est li√©e √† un film",
        "* Une recette peut avoir plusieurs cat√©gories"
    ], image_path="requirements/presentation_assets/mcd.png")
    
    # 15. Base de Donn√©es - MLD
    add_slide_with_content("7.2 Base de Donn√©es - MLD", [
        "**Mod√®le Logique de Donn√©es (MLD)**",
        "",
        "**Tables & cl√©s :**",
        "* UTILISATEUR(id, email, role_id)",
        "* RECETTE(id, film_id, auteur_id)",
        "* FILM(id, titre, annee)",
        "* CATEGORIE(id, libelle)",
        "* RECETTE_CATEGORIE(recette_id, categorie_id)",
        "",
        "**Normalisation :**",
        "* 3NF pour limiter la redondance",
        "* Tables de jonction pour N..N",
        "",
        "**Contraintes :**",
        "* PK sur chaque id",
        "* FK explicites avec r√®gles de suppression"
    ])
    
    # 16. Base de Donn√©es - MPD
    add_slide_with_content("7.3 Base de Donn√©es - MPD", [
        "**Mod√®le Physique de Donn√©es (MPD)**",
        "",
        "**SGBD : PostgreSQL**",
        "",
        "**Types & contraintes :**",
        "* UUID/serial pour identifiants",
        "* VARCHAR pour titres, TEXT pour descriptions",
        "* NOT NULL + UNIQUE sur email",
        "",
        "**Indexation :**",
        "* Index sur film_id, auteur_id",
        "* Index sur recherche de recettes",
        "",
        "**Exploitation :**",
        "* Migrations versionn√©es",
        "* Sauvegardes planifi√©es"
    ])
    
    # 17. Bilan Sprint 0
    add_slide_with_content("8. Bilan Sprint 0", [
        "**Livrables Compl√©t√©s** ‚úÖ",
        "",
        "* ‚úÖ **Cahier des charges** valid√©",
        "* ‚úÖ **Design System** d√©fini",
        "   * Charte graphique",
        "   * Couleurs & Typographie",
        "",
        "* ‚úÖ **Architecture BDD** pr√™te",
        "   * MCD, MLD & MPD valid√©s",
        "",
        "* ‚úÖ **Maquettes** finalis√©es",
        "* ‚úÖ **Wireframes** valid√©s",
        "* ‚úÖ **Arborescence** d√©finie",
        "",
        "**Tous les √©l√©ments sont pr√™ts !**"
    ])
    
    # 18. Prochaines √âtapes
    add_slide_with_content("8.1 Prochaines √âtapes - Sprint 1", [
        "**Cap sur le Sprint 1** üöÄ",
        "",
        "**Objectifs Techniques :**",
        "",
        "* **Initialisation**",
        "   * D√©p√¥t Git & CI/CD",
        "   * Structure du code",
        "",
        "* **Base de donn√©es**",
        "   * Sch√©ma PostgreSQL",
        "   * Migrations & Seeds",
        "",
        "* **Backend (API)**",
        "   * Framework & Routes",
        "   * Authentification",
        "",
        "* **Frontend**",
        "   * Squelette app",
        "   * Design system",
        "   * Premi√®res pages"
    ])
    
    # 19. Conclusion
    add_slide_with_content("Conclusion", [
        "**Cin√© D√©lices - Projet en Cours**",
        "",
        "Une plateforme innovante qui transforme",
        "la passion du cin√©ma en exp√©rience culinaire",
        "",
        "**Merci de votre attention !**",
        "",
        "Questions & Discussion"
    ])
    
    # Save presentation
    output_path = 'requirements/presentation_sprint_0_enhanced.pptx'
    prs.save(output_path)
    print(f"‚úÖ Pr√©sentation am√©lior√©e g√©n√©r√©e : {output_path}")
    print(f"üìä Nombre de slides : {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    create_enhanced_presentation()
